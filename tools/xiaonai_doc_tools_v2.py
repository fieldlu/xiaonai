#!/usr/bin/env python3
"""XiaoNai Document Tools v2 - Universal document reader & creator.

Read any document:
  python3 tools/xiaonai_doc_tools_v2.py read <file> [-n N]

Create documents:
  python3 tools/xiaonai_doc_tools_v2.py make docx <out> <title> <content>
  python3 tools/xiaonai_doc_tools_v2.py make xlsx <out> <json>
  python3 tools/xiaonai_doc_tools_v2.py make pdf  <out> <title> <content>
  python3 tools/xiaonai_doc_tools_v2.py make pptx <out> <title> <content>
  python3 tools/xiaonai_doc_tools_v2.py make md   <out> <content>

Convert:
  python3 tools/xiaonai_doc_tools_v2.py convert <in> <out>
"""
import sys, os, subprocess, json, io


def safe_read(func, *args, **kwargs):
    try:
        return func(*args, **kwargs)
    except Exception as e:
        return f"[Error: {e}]"


def read_txt(path, n=100):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        lines = f.readlines()
    for i, line in enumerate(lines[:n], 1):
        print(f"{i:4d}| {line.rstrip()}")
    if len(lines) > n:
        print(f"... ({len(lines) - n} more lines)")


def read_docx(path, n=60):
    from docx import Document
    doc = Document(path)
    i = 0
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            i += 1
            style = para.style.name if para.style else ""
            prefix = f"[{style}] " if "Heading" in style else ""
            print(f"{i:4d}| {prefix}{t[:200]}")
            if i >= n:
                break
    for table in doc.tables:
        print("--- [TABLE] ---")
        for row in table.rows[:10]:
            cells = [c.text.strip()[:40] for c in row.cells]
            print(" | ".join(c for c in cells if c))


def read_pdf(path, n=100):
    import pdfplumber
    with pdfplumber.open(path) as pdf:
        count = 0
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                for line in text.split("\n"):
                    count += 1
                    print(f"{count:4d}| {line.rstrip()[:200]}")
                    if count >= n:
                        return
            tables = page.extract_tables()
            if tables:
                for t in tables:
                    if t:
                        print("--- [TABLE] ---")
                        for row in t[:5]:
                            print(" | ".join(str(c)[:30] if c else "" for c in row))


def read_pptx(path, n=50):
    from pptx import Presentation
    prs = Presentation(path)
    count = 0
    for slide in prs.slides:
        print(f"\n=== Slide {slide.slide_id} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        count += 1
                        print(f"  {t[:200]}")
                        if count >= n:
                            return
            if shape.has_table:
                print("  [TABLE]")
                for row in shape.table.rows[:5]:
                    print("  | " + " | ".join(c.text.strip()[:20] for c in row.cells))


def read_xlsx(path, n=60):
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    for name in wb.sheetnames:
        ws = wb[name]
        print(f"\n=== {name} ({ws.max_row}r x {ws.max_column}c) ===")
        for row in ws.iter_rows(max_row=min(n, ws.max_row), values_only=True):
            line = " | ".join(str(c)[:60] if c is not None else "" for c in row)
            if line.strip().replace(" | ", ""):
                print(line)
        if ws.max_row > n:
            print(f"... ({ws.max_row - n} more rows)")


def read_doc(path, n=100):
    result = subprocess.run(["antiword", path], capture_output=True, text=True)
    if result.returncode == 0:
        lines = result.stdout.split("\n")
        for i, line in enumerate(lines[:n], 1):
            print(f"{i:4d}| {line.rstrip()}")
    else:
        # Try catdoc as fallback
        result2 = subprocess.run(["catdoc", path], capture_output=True, text=True)
        if result2.returncode == 0:
            for i, line in enumerate(result2.stdout.split("\n")[:n], 1):
                print(f"{i:4d}| {line.rstrip()}")
        else:
            print("[Binary .doc file - use .docx for better support]")


def smart_read(path, n=80):
    """Auto-detect format and read."""
    ext = os.path.splitext(path)[1].lower()
    readers = {
        ".docx": read_docx, ".doc": read_doc, ".pdf": read_pdf,
        ".pptx": read_pptx, ".xlsx": read_xlsx, ".xls": read_xlsx,
        ".csv": read_txt,
        ".txt": read_txt, ".md": read_txt, ".py": read_txt,
        ".json": read_txt, ".html": read_txt, ".xml": read_txt,
    }
    reader = readers.get(ext, read_txt)
    result = safe_read(reader, path, n)
    if isinstance(result, str) and result.startswith("[Error:"):
        print(result)
        return


# ====== CREATE ======

def make_docx(out, title, content):
    from docx import Document
    from docx.shared import Pt, Inches
    doc = Document()
    style = doc.styles["Normal"]
    style.font.size = Pt(11)
    style.font.name = "Arial"
    doc.add_heading(title, 0)
    for line in content.replace("\\n", "\n").split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:], 1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], 2)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(out)
    print(f"Created: {out} ({os.path.getsize(out)} bytes)")


def make_xlsx(out, data_json):
    from openpyxl import Workbook
    from openpyxl.styles import Font
    data = json.loads(data_json)
    wb = Workbook()
    first = True
    for sheet_name, rows in data.items():
        if first:
            ws = wb.active
            ws.title = sheet_name
            first = False
        else:
            ws = wb.create_sheet(sheet_name)
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row, 1):
                cell = ws.cell(row=i, column=j, value=val)
                if i == 1:
                    cell.font = Font(bold=True)
    wb.save(out)
    print(f"Created: {out} ({os.path.getsize(out)} bytes)")


def make_pdf(out, title, content):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    doc = SimpleDocTemplate(out, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for line in content.replace("\\n", "\n").split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            story.append(Paragraph(line[2:], styles["Heading1"]))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles["Heading2"]))
        else:
            story.append(Paragraph(line, styles["Normal"]))
        story.append(Spacer(1, 6))
    doc.build(story)
    print(f"Created: {out} ({os.path.getsize(out)} bytes)")


def make_pptx(out, title, content):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    # Content slides - split by ## headings
    slides_content = []
    current = []
    for line in content.replace("\\n", "\n").split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("## ") and current:
            slides_content.append(current)
            current = [line[3:]]
        else:
            current.append(line)
    if current:
        slides_content.append(current)
    for slide_lines in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide_lines:
            slide.shapes.title.text = slide_lines[0][:80]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, line in enumerate(slide_lines[1:10]):
            if i == 0:
                tf.paragraphs[0].text = line[:200]
            else:
                p = tf.add_paragraph()
                p.text = line[:200]
                p.level = 0 if not line.startswith("- ") else 1
    prs.save(out)
    print(f"Created: {out} ({os.path.getsize(out)} bytes)")


def make_md(out, title, content):
    with open(out, "w", encoding="utf-8") as f:
        f.write(f"# {title}\n\n{content}".replace("\\n", "\n"))
    print(f"Created: {out} ({os.path.getsize(out)} bytes)")


# ====== CONVERT ======

def convert(in_path, out_path):
    ext_in = os.path.splitext(in_path)[1].lower()
    ext_out = os.path.splitext(out_path)[1].lower()

    # Extract content
    import io as _io
    old_stdout = sys.stdout
    import io as _io2; captured = _io2.StringIO(); sys.stdout = captured
    smart_read(in_path, 9999)
    sys.stdout = old_stdout
    content = captured.getvalue()

    # Remove line numbers from read output
    clean = []
    for line in content.split("\n"):
        if "| " in line[:8]:
            clean.append(line.split("| ", 1)[1] if "| " in line else line)
        else:
            clean.append(line)
    content = "\n".join(clean)

    title = os.path.splitext(os.path.basename(in_path))[0]
    if ext_out == ".docx":
        make_docx(out_path, title, content)
    elif ext_out == ".pdf":
        make_pdf(out_path, title, content)
    elif ext_out == ".pptx":
        make_pptx(out_path, title, content)
    elif ext_out in (".xlsx", ".csv"):
        make_xlsx(out_path, json.dumps({"Sheet1": [["Content"], [content]]}))
    else:
        make_md(out_path, content)


# ====== MAIN ======
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("XiaoNai Doc Tools v2")
        print("  read <file> [-n N]     Read document (auto-detect format)")
        print("  make <fmt> <out> ...   Create docx/xlsx/pdf/pptx/md")
        print("  convert <in> <out>     Convert between formats")
        sys.exit(0)

    cmd = sys.argv[1]

    if cmd == "read":
        path = sys.argv[2]
        n = int(sys.argv[4]) if len(sys.argv) > 4 and sys.argv[3] == "-n" else 80
        smart_read(path, n)

    elif cmd == "make":
        fmt, out = sys.argv[2], sys.argv[3]
        rest = " ".join(sys.argv[4:])
        makers = {"docx": make_docx, "xlsx": make_xlsx, "pdf": make_pdf, "pptx": make_pptx, "md": make_md}
        if fmt in makers:
            if fmt == "xlsx":
                makers[fmt](out, rest)
            else:
                title = sys.argv[4] if len(sys.argv) > 4 else "Document"
                content = " ".join(sys.argv[5:]) if len(sys.argv) > 5 else ""
                makers[fmt](out, title, content)
        else:
            print(f"Unknown format: {fmt}")

    elif cmd == "convert":
        convert(sys.argv[2], sys.argv[3])
